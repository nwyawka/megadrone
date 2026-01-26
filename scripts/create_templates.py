#!/usr/bin/env python3
"""
Create JP Security Word and PowerPoint templates with company branding.
"""

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.table import WD_TABLE_ALIGNMENT

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import os

# Paths
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
ASSETS_DIR = os.path.join(PROJECT_DIR, "assets")
TEMPLATES_DIR = os.path.join(ASSETS_DIR, "templates")
LOGO_PATH = os.path.join(ASSETS_DIR, "jp_security_logo.jpg")

# Brand colors
NAVY_BLUE = RGBColor(0x1a, 0x2a, 0x3a)  # Dark navy
GOLD = RGBColor(0xc9, 0xa2, 0x5c)  # Gold accent
CREAM = RGBColor(0xf5, 0xf0, 0xe1)  # Cream/off-white

def create_word_template():
    """Create JP Security Word document template."""
    doc = Document()

    # Set up styles
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)
    font.color.rgb = NAVY_BLUE

    # Title style
    title_style = doc.styles['Title']
    title_style.font.name = 'Calibri'
    title_style.font.size = Pt(24)
    title_style.font.bold = True
    title_style.font.color.rgb = NAVY_BLUE

    # Heading 1 style
    h1_style = doc.styles['Heading 1']
    h1_style.font.name = 'Calibri'
    h1_style.font.size = Pt(16)
    h1_style.font.bold = True
    h1_style.font.color.rgb = NAVY_BLUE

    # Heading 2 style
    h2_style = doc.styles['Heading 2']
    h2_style.font.name = 'Calibri'
    h2_style.font.size = Pt(14)
    h2_style.font.bold = True
    h2_style.font.color.rgb = NAVY_BLUE

    # Add header with logo
    section = doc.sections[0]
    header = section.header
    header_para = header.paragraphs[0]
    header_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Add logo to header
    run = header_para.add_run()
    run.add_picture(LOGO_PATH, width=Inches(1.5))

    # Add company name below logo
    header_para2 = header.add_paragraph()
    header_para2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run2 = header_para2.add_run("JP SECURITY")
    run2.font.size = Pt(14)
    run2.font.bold = True
    run2.font.color.rgb = NAVY_BLUE

    # Add tagline
    header_para3 = header.add_paragraph()
    header_para3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run3 = header_para3.add_run("FORTIS ET FIDELIS")
    run3.font.size = Pt(10)
    run3.font.italic = True
    run3.font.color.rgb = GOLD

    # Add horizontal line
    header_para4 = header.add_paragraph()
    header_para4.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run4 = header_para4.add_run("─" * 60)
    run4.font.size = Pt(8)
    run4.font.color.rgb = GOLD

    # Add footer
    footer = section.footer
    footer_para = footer.paragraphs[0]
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Footer line
    footer_run = footer_para.add_run("─" * 60)
    footer_run.font.size = Pt(8)
    footer_run.font.color.rgb = GOLD

    footer_para2 = footer.add_paragraph()
    footer_para2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run2 = footer_para2.add_run("JP Security  |  Maritime ISR & Privateer Marque  |  UNCLASSIFIED // PROPRIETARY")
    footer_run2.font.size = Pt(9)
    footer_run2.font.color.rgb = NAVY_BLUE

    # Add document body placeholder content
    doc.add_paragraph()
    doc.add_paragraph()

    # Title placeholder
    title = doc.add_paragraph("[DOCUMENT TITLE]", style='Title')
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Metadata table
    doc.add_paragraph()
    table = doc.add_table(rows=4, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    metadata = [
        ("Classification:", "UNCLASSIFIED // PROPRIETARY"),
        ("Version:", "[X.X]"),
        ("Date:", "[Date]"),
        ("Author:", "[Author Name]"),
    ]

    for i, (label, value) in enumerate(metadata):
        row = table.rows[i]
        row.cells[0].text = label
        row.cells[1].text = value
        row.cells[0].paragraphs[0].runs[0].font.bold = True

    doc.add_paragraph()
    doc.add_paragraph("─" * 40)
    doc.add_paragraph()

    # Section placeholders
    doc.add_heading("1. Section Heading", level=1)
    doc.add_paragraph("[Section content goes here...]")
    doc.add_paragraph()

    doc.add_heading("1.1 Subsection Heading", level=2)
    doc.add_paragraph("[Subsection content goes here...]")
    doc.add_paragraph()

    doc.add_heading("2. Section Heading", level=1)
    doc.add_paragraph("[Section content goes here...]")

    # Save template
    os.makedirs(TEMPLATES_DIR, exist_ok=True)
    output_path = os.path.join(TEMPLATES_DIR, "JP_Security_Template.docx")
    doc.save(output_path)
    print(f"Word template saved to: {output_path}")
    return output_path


def create_pptx_template():
    """Create JP Security PowerPoint template."""
    prs = Presentation()

    # Set slide dimensions (16:9)
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add background color (navy)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PptxRGBColor(0x1a, 0x2a, 0x3a)

    # Add logo (centered, large)
    logo = slide.shapes.add_picture(
        LOGO_PATH,
        PptxInches(4.5),
        PptxInches(1.0),
        width=PptxInches(4.333)
    )

    # Add title text box
    title_box = slide.shapes.add_textbox(
        PptxInches(0.5),
        PptxInches(5.5),
        PptxInches(12.333),
        PptxInches(1.0)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "[PRESENTATION TITLE]"
    title_para.font.size = PptxPt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PptxRGBColor(0xf5, 0xf0, 0xe1)  # Cream
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        PptxInches(0.5),
        PptxInches(6.3),
        PptxInches(12.333),
        PptxInches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "[Date] | [Presenter Name]"
    subtitle_para.font.size = PptxPt(18)
    subtitle_para.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)  # Gold
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Classification footer
    class_box = slide.shapes.add_textbox(
        PptxInches(0.5),
        PptxInches(7.0),
        PptxInches(12.333),
        PptxInches(0.3)
    )
    class_frame = class_box.text_frame
    class_para = class_frame.paragraphs[0]
    class_para.text = "UNCLASSIFIED // PROPRIETARY"
    class_para.font.size = PptxPt(10)
    class_para.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)  # Gold
    class_para.alignment = PP_ALIGN.CENTER

    # Content slide
    slide2 = prs.slides.add_slide(slide_layout)

    # Add background
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = PptxRGBColor(0x1a, 0x2a, 0x3a)

    # Add small logo in corner
    logo2 = slide2.shapes.add_picture(
        LOGO_PATH,
        PptxInches(0.3),
        PptxInches(0.2),
        width=PptxInches(0.8)
    )

    # Add slide title
    title_box2 = slide2.shapes.add_textbox(
        PptxInches(1.3),
        PptxInches(0.3),
        PptxInches(11.5),
        PptxInches(0.7)
    )
    title_frame2 = title_box2.text_frame
    title_para2 = title_frame2.paragraphs[0]
    title_para2.text = "[Slide Title]"
    title_para2.font.size = PptxPt(32)
    title_para2.font.bold = True
    title_para2.font.color.rgb = PptxRGBColor(0xf5, 0xf0, 0xe1)  # Cream

    # Add horizontal line
    line = slide2.shapes.add_shape(
        1,  # Rectangle
        PptxInches(0.5),
        PptxInches(1.1),
        PptxInches(12.333),
        PptxInches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)  # Gold
    line.line.fill.background()

    # Add content placeholder
    content_box = slide2.shapes.add_textbox(
        PptxInches(0.5),
        PptxInches(1.4),
        PptxInches(12.333),
        PptxInches(5.3)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_para = content_frame.paragraphs[0]
    content_para.text = "[Slide content goes here...]"
    content_para.font.size = PptxPt(20)
    content_para.font.color.rgb = PptxRGBColor(0xf5, 0xf0, 0xe1)  # Cream

    # Add bullet points
    p2 = content_frame.add_paragraph()
    p2.text = "• First bullet point"
    p2.font.size = PptxPt(18)
    p2.font.color.rgb = PptxRGBColor(0xf5, 0xf0, 0xe1)
    p2.level = 0

    p3 = content_frame.add_paragraph()
    p3.text = "• Second bullet point"
    p3.font.size = PptxPt(18)
    p3.font.color.rgb = PptxRGBColor(0xf5, 0xf0, 0xe1)
    p3.level = 0

    p4 = content_frame.add_paragraph()
    p4.text = "    - Sub-bullet point"
    p4.font.size = PptxPt(16)
    p4.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)  # Gold for sub-bullets
    p4.level = 1

    # Footer
    footer_box = slide2.shapes.add_textbox(
        PptxInches(0.5),
        PptxInches(7.0),
        PptxInches(6),
        PptxInches(0.3)
    )
    footer_frame = footer_box.text_frame
    footer_para = footer_frame.paragraphs[0]
    footer_para.text = "JP Security | Fortis et Fidelis"
    footer_para.font.size = PptxPt(10)
    footer_para.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)

    # Page number placeholder
    page_box = slide2.shapes.add_textbox(
        PptxInches(12.0),
        PptxInches(7.0),
        PptxInches(1),
        PptxInches(0.3)
    )
    page_frame = page_box.text_frame
    page_para = page_frame.paragraphs[0]
    page_para.text = "2"
    page_para.font.size = PptxPt(10)
    page_para.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)
    page_para.alignment = PP_ALIGN.RIGHT

    # Classification footer
    class_box2 = slide2.shapes.add_textbox(
        PptxInches(4.0),
        PptxInches(7.0),
        PptxInches(5.333),
        PptxInches(0.3)
    )
    class_frame2 = class_box2.text_frame
    class_para2 = class_frame2.paragraphs[0]
    class_para2.text = "UNCLASSIFIED // PROPRIETARY"
    class_para2.font.size = PptxPt(10)
    class_para2.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)
    class_para2.alignment = PP_ALIGN.CENTER

    # Two-column slide
    slide3 = prs.slides.add_slide(slide_layout)

    # Background
    background3 = slide3.background
    fill3 = background3.fill
    fill3.solid()
    fill3.fore_color.rgb = PptxRGBColor(0x1a, 0x2a, 0x3a)

    # Logo
    logo3 = slide3.shapes.add_picture(
        LOGO_PATH,
        PptxInches(0.3),
        PptxInches(0.2),
        width=PptxInches(0.8)
    )

    # Title
    title_box3 = slide3.shapes.add_textbox(
        PptxInches(1.3),
        PptxInches(0.3),
        PptxInches(11.5),
        PptxInches(0.7)
    )
    title_frame3 = title_box3.text_frame
    title_para3 = title_frame3.paragraphs[0]
    title_para3.text = "[Two-Column Layout]"
    title_para3.font.size = PptxPt(32)
    title_para3.font.bold = True
    title_para3.font.color.rgb = PptxRGBColor(0xf5, 0xf0, 0xe1)

    # Horizontal line
    line3 = slide3.shapes.add_shape(
        1,
        PptxInches(0.5),
        PptxInches(1.1),
        PptxInches(12.333),
        PptxInches(0.02)
    )
    line3.fill.solid()
    line3.fill.fore_color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)
    line3.line.fill.background()

    # Left column
    left_box = slide3.shapes.add_textbox(
        PptxInches(0.5),
        PptxInches(1.4),
        PptxInches(5.9),
        PptxInches(5.3)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_para = left_frame.paragraphs[0]
    left_para.text = "Left Column"
    left_para.font.size = PptxPt(20)
    left_para.font.bold = True
    left_para.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)

    left_p2 = left_frame.add_paragraph()
    left_p2.text = "[Content for left column...]"
    left_p2.font.size = PptxPt(16)
    left_p2.font.color.rgb = PptxRGBColor(0xf5, 0xf0, 0xe1)

    # Right column
    right_box = slide3.shapes.add_textbox(
        PptxInches(6.9),
        PptxInches(1.4),
        PptxInches(5.9),
        PptxInches(5.3)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_para = right_frame.paragraphs[0]
    right_para.text = "Right Column"
    right_para.font.size = PptxPt(20)
    right_para.font.bold = True
    right_para.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)

    right_p2 = right_frame.add_paragraph()
    right_p2.text = "[Content for right column...]"
    right_p2.font.size = PptxPt(16)
    right_p2.font.color.rgb = PptxRGBColor(0xf5, 0xf0, 0xe1)

    # Footer for slide 3
    footer_box3 = slide3.shapes.add_textbox(
        PptxInches(0.5),
        PptxInches(7.0),
        PptxInches(6),
        PptxInches(0.3)
    )
    footer_frame3 = footer_box3.text_frame
    footer_para3 = footer_frame3.paragraphs[0]
    footer_para3.text = "JP Security | Fortis et Fidelis"
    footer_para3.font.size = PptxPt(10)
    footer_para3.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)

    class_box3 = slide3.shapes.add_textbox(
        PptxInches(4.0),
        PptxInches(7.0),
        PptxInches(5.333),
        PptxInches(0.3)
    )
    class_frame3 = class_box3.text_frame
    class_para3 = class_frame3.paragraphs[0]
    class_para3.text = "UNCLASSIFIED // PROPRIETARY"
    class_para3.font.size = PptxPt(10)
    class_para3.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)
    class_para3.alignment = PP_ALIGN.CENTER

    page_box3 = slide3.shapes.add_textbox(
        PptxInches(12.0),
        PptxInches(7.0),
        PptxInches(1),
        PptxInches(0.3)
    )
    page_frame3 = page_box3.text_frame
    page_para3 = page_frame3.paragraphs[0]
    page_para3.text = "3"
    page_para3.font.size = PptxPt(10)
    page_para3.font.color.rgb = PptxRGBColor(0xc9, 0xa2, 0x5c)
    page_para3.alignment = PP_ALIGN.RIGHT

    # Save template
    os.makedirs(TEMPLATES_DIR, exist_ok=True)
    output_path = os.path.join(TEMPLATES_DIR, "JP_Security_Template.pptx")
    prs.save(output_path)
    print(f"PowerPoint template saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    print("Creating JP Security templates...")
    print(f"Using logo: {LOGO_PATH}")
    print()

    word_path = create_word_template()
    pptx_path = create_pptx_template()

    print()
    print("Templates created successfully!")
    print(f"Word template: {word_path}")
    print(f"PowerPoint template: {pptx_path}")
